from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os,sys
 

def add_slide_with_picture(pptx_path, img_path, title='title'):
    im = Image.open(img_path)
    iw, ih = im.size
    
    # load a presentation and get its dimensions
    prs = Presentation(pptx_path)
    sh, sw = (prs.slide_height, prs.slide_width)
    #print(sh, sw)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    ## Add image fitting the width of the slide
    top = sh-sw*ih/iw # align image at the bottom
    #top = 0.5*(sh-sw*ih/iw) # align image at the center
    left = 0. 
    pic = slide.shapes.add_picture(img_path, left, top, width=sw)
    
    ## Add text box
    top = 0.
    left = 0.
    txBox = slide.shapes.add_textbox(left, top, width=sw, height=sh*0.05)
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(pptx_path)

def add_slide_recursively():
    pptx_path = 'test_python_pptx.pptx'
    #img_path = '/data/c3s_vol6/TEST_CNRM/remymf_test/vegeo_trend_analysis/tmp/mosaic_YYYY.png'
    img_path = '/data/c3s_vol6/TEST_CNRM/remymf_test/vegeo_trend_analysis/tmp/mosaic_snht_YYYY.png'
    res_pptx_path = 'out.pptx'
    os.system('cp {} {}'.format(pptx_path, res_pptx_path))

    #for year in [str(i) for i in range(1982,2006)]:
    for year in ['{:02d}'.format(i) for i in range(1,16)]:
        print(year)
        add_slide_with_picture(res_pptx_path, img_path.replace('YYYY',year), year)


if __name__=='__main__':

    #add_slide_with_picture()
    add_slide_recursively()
